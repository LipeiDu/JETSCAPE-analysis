"""
  macro for plotting analyzed jetscape events
  """

# This script plots histograms created in the analysis of Jetscape events
#
# Author: James Mulligan (james.mulligan@berkeley.edu)

# General
import ctypes
import os
import sys
import yaml
import argparse

# Data analysis and plotting
import ROOT
import numpy as np
import uproot
import pandas as pd
import h5py

# Base class
sys.path.append('.')
from jetscape_analysis.base import common_base
from plot import plot_results_STAT_utils

# Prevent ROOT from stealing focus when plotting
ROOT.gROOT.SetBatch(True)

################################################################
class PlotResults(common_base.CommonBase):

    # ---------------------------------------------------------------
    # Constructor
    # ---------------------------------------------------------------
    def __init__(self, config_file='', input_file='', output_dir='', **kwargs):
        super(PlotResults, self).__init__(**kwargs)

        if output_dir:
            self.output_dir = output_dir
            if not os.path.exists(self.output_dir):
                os.makedirs(self.output_dir)
        else:
            self.output_dir = os.path.dirname(input_file)

        self.plot_utils = plot_results_STAT_utils.PlotUtils()
        self.plot_utils.setOptions()
        ROOT.gROOT.ForceStyle()

        self.input_file = ROOT.TFile(input_file, 'READ')

        self.data_color = ROOT.kGray+3
        self.data_marker = 21
        self.jetscape_color = [ROOT.kViolet-8, ROOT.kViolet-8, ROOT.kRed-7, ROOT.kTeal-8, ROOT.kCyan-2, ROOT.kGreen-6, ROOT.kAzure-4, ROOT.kOrange+6, ROOT.kBlue-10]
        self.jetscape_fillstyle = [1001, 3144, 1001, 3144]
        self.jetscape_alpha = [0.7, 0.7, 0.7, 0.7]
        self.jetscape_marker = 20
        self.marker_size = 1.5
        self.line_width = 2
        self.line_style = 1
        self.file_format = '.pdf'

        # Check whether pp or AA
        self.is_AA = True
        self.observable_centrality_list = []

        # Read config file
        with open(config_file, 'r') as stream:
            self.config = yaml.safe_load(stream)

        self.sqrts = self.config['sqrt_s']
        self.n_order = self.config['n_order']

        # for pT interpolation in flow analysis
        self.pT_low = self.config['pT_low']
        self.pT_high = self.config['pT_high']
        self.npT = self.config['npT']

        # Qn vector parameters
        self.n_pt_bins = self.config['n_pt_bins']
        self.pt_min = self.config['pt_min']
        self.pt_max = self.config['pt_max']

        self.pt_array = np.linspace(self.pt_min, self.pt_max, self.n_pt_bins)

        # We will write final results after all scalings, along with data, to file
        self.output_dict = {}

        print(self)

    #-------------------------------------------------------------------------------------------
    #-------------------------------------------------------------------------------------------
    #-------------------------------------------------------------------------------------------
    def plot_results(self):

        self.plot_soft_observable(observable_type='soft_integrated')

        self.plot_soft_observable(observable_type='soft_differential')
        
        # self.plot_event_qa()

        # self.write_output_objects()

        # # Generate pptx for convenience
        # if self.file_format == '.png':
        #     self.generate_pptx()

    #-------------------------------------------------------------------------------------------
    # Plot hadron observables
    #-------------------------------------------------------------------------------------------
    def plot_soft_observable(self, observable_type=''):
        print()
        print(f'Plot {observable_type} observables...')

        for observable, block in self.config[observable_type].items():
            for centrality_index,centrality in enumerate(block['centrality']):

                if 'hepdata' not in block and 'custom_data' not in block and 'user_data' not in block:
                    continue

                # Initialize observable configuration
                self.suffix = ''
                self.init_observable(observable_type, observable, block, centrality, centrality_index)

                # Plot observable
                self.plot_observable(observable_type, observable, centrality)

    #-------------------------------------------------------------------------------------------
    # Initialize a single observable's config
    #-------------------------------------------------------------------------------------------
    def init_observable(self, observable_type, observable, block, centrality, centrality_index, method='', pt_suffix='', self_normalize=False):

        # Initialize an empty dict containing relevant info
        self.observable_settings = {}

        # Initialize common settings into class members
        self.init_common_settings(observable_type, observable, block)

        #-----------------------------------------------------------
        # Initialize data distribution into self.observable_settings
        if 'hepdata' in block:
            self.observable_settings['data_distribution'] = self.plot_utils.tgraph_from_hepdata(block, self.is_AA, self.sqrts, observable_type, observable, centrality_index, suffix=self.suffix, pt_suffix=pt_suffix)
        elif 'custom_data' in block:
            self.observable_settings['data_distribution'] = self.plot_utils.tgraph_from_yaml(block, self.is_AA, self.sqrts, observable_type, observable, centrality_index, suffix=self.suffix, pt_suffix=pt_suffix)
        elif 'user_data' in block:
            self.observable_settings['data_distribution'] = self.plot_utils.tgraph_from_custom_yaml(block, self.sqrts, observable_type, observable, centrality_index = None)
        else:
            self.observable_settings['data_distribution'] = None

        #-----------------------------------------------------------
        # Initialize JETSCAPE distribution into self.observable_settings
        self.initialize_jetscape_distribution(observable_type, observable, method, block, centrality, centrality_index, pt_suffix=pt_suffix, self_normalize=self_normalize)

        # #-----------------------------------------------------------
        # # For v2 of hadrons -- form ratio of JETSCAPE to data and load into self.observable_settings
        # if observable_type in ['hadron_correlations'] and 'v2' in observable and self.is_AA and self.observable_settings[f'jetscape_distribution']:
        #     self.observable_settings['ratio'] =  self.plot_utils.divide_histogram_by_tgraph(self.observable_settings[f'jetscape_distribution'], self.observable_settings['data_distribution'])

        # #-----------------------------------------------------------
        # # For v2 of jets -- form ratio of JETSCAPE to data and load into self.observable_settings; currently observable_settings['ratio'] is overwritten, need fix; LDU
        # if observable_type in ['dijet'] and 'v2' in observable and self.is_AA:
        #     for jet_collection_label in self.jet_collection_labels:
        #         if self.observable_settings[f'jetscape_distribution{jet_collection_label}']:
        #             self.observable_settings['ratio'] =  self.plot_utils.divide_histogram_by_tgraph(self.observable_settings[f'jetscape_distribution{jet_collection_label}'], self.observable_settings['data_distribution'])

    #-------------------------------------------------------------------------------------------
    # Initialize from settings from config file into class members
    #-------------------------------------------------------------------------------------------
    def init_common_settings(self, observable_type, observable, block):

        self.xtitle = block['xtitle']
        if 'eta_cut' in block:
            self.eta_cut = block['eta_cut']
        if 'y_cut' in block:
            self.y_cut = block['y_cut']
        if 'pt' in block:
            self.pt = block['pt']
        if 'eta_cut_R' in block:
            self.eta_R = block['eta_cut_R']
            self.eta_cut = np.round(self.eta_R - self.jet_R, decimals=1)
        if 'c_ref' in block:
            index = block['jet_R'].index(self.jet_R)
            self.c_ref = block['c_ref'][index]
        if 'low_trigger_range' in block:
            self.low_trigger_range = block['low_trigger_range']
        if 'high_trigger_range' in block:
            self.high_trigger_range = block['high_trigger_range']
        if 'trigger_range' in block:
            self.trigger_range = block['trigger_range']
        if 'logy' in block:
            self.logy = block['logy']
        else:
            self.logy = False

        if self.is_AA:
            if 'ytitle_AA' in block:
                self.ytitle = block['ytitle_AA']
            if 'y_min_AA' in block:
                self.y_min = float(block['y_min_AA'])
                self.y_max = float(block['y_max_AA'])
            else:
                self.y_min = 0.
                self.y_max = 1.
        else:
            if 'ytitle_pp' in block:
                self.ytitle = block['ytitle_pp']
            else:
                self.ytitle = ''
            if 'y_min_pp' in block:
                self.y_min = float(block['y_min_pp'])
                self.y_max = float(block['y_max_pp'])
            else:
                self.y_min = 0.
                self.y_max = 1.99
            if 'y_ratio_min' in block:
                self.y_ratio_min = block['y_ratio_min']
                self.y_ratio_max = block['y_ratio_max']
            else:
                self.y_ratio_min = 0.
                self.y_ratio_max = 1.99
            # Provide the opportunity to disable logy in pp, but keep in AA
            if 'logy_pp' in block:
                self.logy = block["logy_pp"]

        #for v2
        if 'integrated' in observable_type:
            if 'v2' in observable:
                self.y_min = -0.05
                self.y_max = 0.15
                self.y_ratio_min = -0.5
                self.y_ratio_max = 1.99
        if 'differential' in observable_type:
            if 'v2' in observable:
                self.y_min = -0.05
                self.y_max = 0.5
                self.y_ratio_min = -0.5
                self.y_ratio_max = 1.99


        # for multiplicity
        if 'multiplicity' in observable:
            self.y_max = 500

        if 'vn' in observable:
            self.y_max = 0.2

        if 'skip_pp' in block:
            self.skip_pp = block['skip_pp']
        else:
            self.skip_pp = False
        if 'skip_pp_ratio' in block:
            self.skip_pp_ratio = block['skip_pp_ratio']
        else:
            self.skip_pp_ratio = False
        if 'skip_AA_ratio' in block:
            self.skip_AA_ratio = block['skip_AA_ratio']
        else:
            self.skip_AA_ratio = False
        if 'scale_by' in block:
            self.scale_by = block['scale_by']
        else:
            self.scale_by = None

        # Flag to plot hole histogram (for hadron histograms only)
        if self.is_AA:
            self.subtract_holes = observable_type in ['hadron', 'hadron_correlations'] and observable in ['pt_ch_alice', 'pt_pi_alice', 'pt_pi0_alice', 'pt_ch_cms',
                                                 'pt_ch_atlas', 'pt_pi0_phenix', 'pt_ch_star', 'v2_atlas', 'v2_cms', 'v2_alice',
                                                 'dihadron_star']
        else:
            self.subtract_holes = False

    #-------------------------------------------------------------------------------------------
    # Initialize JETSCAPE distribution into self.observable_settings
    #-------------------------------------------------------------------------------------------
    def initialize_jetscape_distribution(self, observable_type, observable, method, block, centrality, centrality_index, pt_suffix='', self_normalize=False):

        # Add centrality bin to list, if needed
        if centrality not in self.observable_centrality_list:
            self.observable_centrality_list.append(centrality)

        self.get_histogram(observable_type, observable, centrality)
        self.post_process_histogram(observable_type, observable, block, centrality)

    #-------------------------------------------------------------------------------------------
    # Get histogram and add to self.observable_settings
    #-------------------------------------------------------------------------------------------
    def get_histogram(self, observable_type, observable, centrality):
        keys = [key.GetName() for key in self.input_file.GetListOfKeys()]

        # Common histogram name prefix
        base_name = f"h_{observable_type}_{observable}_{centrality}"

        if observable_type == "soft_integrated":
            if 'multiplicity' in observable:
                hist_name = f"{base_name}_dNchdeta"
                if hist_name in keys:
                    self.observable_settings['hist_dNchdeta'] = self.input_file.Get(hist_name)
            
            if 'v2' in observable:
                self.observable_settings['hist_N_vn'] = self.input_file.Get(f"{base_name}_N_vn")

                self.observable_settings['hist_vn_real'] = {}
                self.observable_settings['hist_vn_imag'] = {}

                for n in range(1, self.n_order):
                    real_name = f"{base_name}_vn_real_n{n}"
                    imag_name = f"{base_name}_vn_imag_n{n}"
                    self.observable_settings['hist_vn_real'][n] = self.input_file.Get(real_name)
                    self.observable_settings['hist_vn_imag'][n] = self.input_file.Get(imag_name)

        if observable_type == "soft_differential":
            if 'v2' in observable:
                self.observable_settings['hist_N_Qn_pT'] = self.input_file.Get(f"{base_name}_N_Qn_pT")
                self.observable_settings['hist_N_Qn_ref'] = self.input_file.Get(f"{base_name}_N_Qn_ref")

                self.observable_settings['hist_Qn_pT_real'] = {}
                self.observable_settings['hist_Qn_pT_imag'] = {}
                self.observable_settings['hist_Qn_ref_real'] = {}
                self.observable_settings['hist_Qn_ref_imag'] = {}

                for n in range(1, self.n_order):
                    real_name_pT = f"{base_name}_Qn_pT_real_n{n}"
                    imag_name_pT = f"{base_name}_Qn_pT_imag_n{n}"
                    real_name_ref = f"{base_name}_Qn_ref_real_n{n}"
                    imag_name_ref = f"{base_name}_Qn_ref_imag_n{n}"
                    self.observable_settings['hist_Qn_pT_real'][n] = self.input_file.Get(real_name_pT)
                    self.observable_settings['hist_Qn_pT_imag'][n] = self.input_file.Get(imag_name_pT)
                    self.observable_settings['hist_Qn_ref_real'][n] = self.input_file.Get(real_name_ref)
                    self.observable_settings['hist_Qn_ref_imag'][n] = self.input_file.Get(imag_name_ref)

    #-------------------------------------------------------------------------------------------
    # Perform any additional manipulations on scaled histograms
    #-------------------------------------------------------------------------------------------
    def post_process_histogram(self, observable_type, observable, block, centrality):

        base_name = f"jetscape_distribution_{observable_type}_{observable}_{centrality}"

        if observable_type == "soft_integrated":
            if 'multiplicity' in observable:
                hist_dNchdeta = self.observable_settings['hist_dNchdeta']

                if hist_dNchdeta:
                    n_events = hist_dNchdeta.GetNbinsX()
                    tot_dNchdeta = sum(hist_dNchdeta.GetBinContent(i) for i in range(1, n_events + 1))
                    mean_dNchdeta = tot_dNchdeta / n_events if n_events > 0 else 0

                    self.observable_settings[f'{base_name}'] = mean_dNchdeta

            elif 'v2' in observable:
                hist_N_vn = self.observable_settings['hist_N_vn']
                hist_vn_real = self.observable_settings['hist_vn_real']
                hist_vn_imag = self.observable_settings['hist_vn_imag']

                # read histogram and pack results in the right format
                vn_data_array = []
                for event_id in range(1, hist_N_vn.GetNbinsX() + 1):
                    N_event = hist_N_vn.GetBinContent(event_id)
                    temp_vn_array = [N_event]
                    for n in range(1, self.n_order):
                        vn_real_event = hist_vn_real[n].GetBinContent(event_id)
                        vn_imag_event = hist_vn_imag[n].GetBinContent(event_id)
                        temp_vn_array.append(vn_real_event + 1j * vn_imag_event)
                    vn_data_array.append(temp_vn_array)

                # integrated flow calculation
                if "two" in observable:
                    vn2 = self.calculate_vn2(vn_data_array)
                    self.observable_settings[f'{base_name}'] = vn2

                if "four" in observable:
                    vn4 = self.calculate_vn4(vn_data_array)
                    self.observable_settings[f'{base_name}'] = vn4

        if observable_type == "soft_differential":
            if 'v2' in observable:
                hist_N_Qn_pT = self.observable_settings['hist_N_Qn_pT']
                hist_N_Qn_ref = self.observable_settings['hist_N_Qn_ref']
                hist_Qn_pT_real = self.observable_settings['hist_Qn_pT_real']
                hist_Qn_pT_imag = self.observable_settings['hist_Qn_pT_imag']
                hist_Qn_ref_real = self.observable_settings['hist_Qn_ref_real']
                hist_Qn_ref_imag = self.observable_settings['hist_Qn_ref_imag']

                # read histogram and pack results in the right format
                QnpT_diff_array, Qnref_array = [], []

                for event_id in range(1, hist_N_Qn_pT.GetNbinsX() + 1):
                    N_Qn_pT = np.array([hist_N_Qn_pT.GetBinContent(event_id, pT_bin) for pT_bin in range(1, hist_N_Qn_pT.GetNbinsY() + 1)])
                    N_Qn_ref = hist_N_Qn_ref.GetBinContent(event_id)

                    QnpT_diff_event = [N_Qn_pT]
                    Qnref_event = [N_Qn_ref]
                    for n in range(1, self.n_order):
                        Qn_pT_real_n = np.array([hist_Qn_pT_real[n].GetBinContent(event_id, pT_bin) for pT_bin in range(1, hist_N_Qn_pT.GetNbinsY() + 1)])
                        Qn_pT_imag_n = np.array([hist_Qn_pT_imag[n].GetBinContent(event_id, pT_bin) for pT_bin in range(1, hist_N_Qn_pT.GetNbinsY() + 1)])
                        Qn_ref_real = hist_Qn_ref_real[n].GetBinContent(event_id)
                        Qn_ref_imag = hist_Qn_ref_imag[n].GetBinContent(event_id)
                        QnpT_diff_event.append(Qn_pT_real_n + 1j * Qn_pT_imag_n)
                        Qnref_event.append(Qn_ref_real + 1j * Qn_ref_imag)
                    QnpT_diff_array.append(QnpT_diff_event)
                    Qnref_array.append(Qnref_event)

                # integrated flow calculation
                if "SP" in observable:
                    vn_SP_diff = self.calculate_vnSP_diff(QnpT_diff_array, Qnref_array)
                    vn_SP_diff_formatted = {
                        'pT_bins': self.pt_array,  # Using self.pt_array defined in the constructor
                        'vn_values': vn_SP_diff[0],  # vn_values for the different harmonic orders
                        'vn_errors': vn_SP_diff[1]   # vn_errors for the different harmonic orders
                    }
                    self.observable_settings[f'{base_name}'] = vn_SP_diff_formatted

                if "four" in observable:
                    vn4_diff = self.calculate_vn4_diff(QnpT_diff_array, Qnref_array)
                    vn4_diff_formatted = {
                        'pT_bins': self.pt_array,  # Using self.pt_array defined in the constructor
                        'vn_values': vn4_diff[0],  # vn_values for the different harmonic orders
                        'vn_errors': vn4_diff[1]   # vn_errors for the different harmonic orders
                    }
                    self.observable_settings[f'{base_name}'] = vn4_diff_formatted

    #-------------------------------------------------------------------------------------------
    # Functions for flow calculations
    #-------------------------------------------------------------------------------------------
    def calculate_vn2(self, vn_data_array):
        """
            this function computes vn{2} and its stat. err.
            self correlation is substracted
        """
        vn_data_array = np.array(vn_data_array)
        nev = len(vn_data_array[:, 0])
        dN = np.real(vn_data_array[:, 0])
        dN = dN.reshape(len(dN), 1)
        Qn_array = dN*vn_data_array[:, 1:]
        corr = 1./(dN*(dN - 1.))*(Qn_array*np.conj(Qn_array) - dN)
        vn_2 = np.sqrt(np.real(np.mean(corr, 0))) + 1e-30
        vn_2_err = np.std(np.real(corr), 0)/np.sqrt(nev)/2./vn_2
        return (np.nan_to_num(vn_2), np.nan_to_num(vn_2_err))

    def calculate_vn4(self, vn_data_array):
        """
        This function computes the 4-particle cumulant vn{4}
        vn{4} = (2 <v_n*conj(v_n)>**2 - <(v_n*conj(v_n))**2.>)**(1/4)
        """
        vn_data_array = np.array(vn_data_array)
        nev = len(vn_data_array[:, 0])
        dN = np.real(vn_data_array[:, 0])
        dN = dN + 1.e-30

        Q1 = dN * vn_data_array[:, 1]
        Q2 = dN * vn_data_array[:, 2]
        Q3 = dN * vn_data_array[:, 3]
        Q4 = dN * vn_data_array[:, 4]
        Q5 = dN * vn_data_array[:, 5]
        Q6 = dN * vn_data_array[:, 6]

        # Two-particle correlation
        N2_weight = dN * (dN - 1.0)
        Q1_2 = np.abs(Q1) ** 2 - dN
        Q2_2 = np.abs(Q2) ** 2 - dN
        Q3_2 = np.abs(Q3) ** 2 - dN

        # Four-particle correlation
        N4_weight = dN * (dN - 1.0) * (dN - 2.0) * (dN - 3.0) + 1.e-30
        Q1_4 = (np.abs(Q1) ** 4
                - 2.0 * np.real(Q2 * np.conj(Q1) * np.conj(Q1))
                - 4.0 * (dN - 2.0) * np.abs(Q1) ** 2
                + np.abs(Q2) ** 2
                + 2.0 * dN * (dN - 3.0))

        Q2_4 = (np.abs(Q2) ** 4
                - 2.0 * np.real(Q4 * np.conj(Q2) * np.conj(Q2))
                - 4.0 * (dN - 2.0) * np.abs(Q2) ** 2
                + np.abs(Q4) ** 2
                + 2.0 * dN * (dN - 3.0))

        Q3_4 = (np.abs(Q3) ** 4
                - 2.0 * np.real(Q6 * np.conj(Q3) * np.conj(Q3))
                - 4.0 * (dN - 2.0) * np.abs(Q3) ** 2
                + np.abs(Q6) ** 2
                + 2.0 * dN * (dN - 3.0))

        # Calculate observables with Jackknife resampling method
        C1_4_array = np.zeros(nev)
        C2_4_array = np.zeros(nev)
        C3_4_array = np.zeros(nev)
        for iev in range(nev):
            array_idx = np.ones(nev, dtype=bool)
            array_idx[iev] = False

            # C_1{4}
            C1_4_array[iev] = (np.mean(Q1_4[array_idx]) / np.mean(N4_weight[array_idx])
                               - 2.0 * (np.mean(Q1_2[array_idx]) / np.mean(N2_weight[array_idx])) ** 2)
            # C_2{4}
            C2_4_array[iev] = (np.mean(Q2_4[array_idx]) / np.mean(N4_weight[array_idx])
                               - 2.0 * (np.mean(Q2_2[array_idx]) / np.mean(N2_weight[array_idx])) ** 2)
            # C_3{4}
            C3_4_array[iev] = (np.mean(Q3_4[array_idx]) / np.mean(N4_weight[array_idx])
                               - 2.0 * (np.mean(Q3_2[array_idx]) / np.mean(N2_weight[array_idx])) ** 2)

        C1_4_mean = np.mean(C1_4_array)
        C1_4_err = np.sqrt((nev - 1.0) / nev * np.sum((C1_4_array - C1_4_mean) ** 2))
        C2_4_mean = np.mean(C2_4_array)
        C2_4_err = np.sqrt((nev - 1.0) / nev * np.sum((C2_4_array - C2_4_mean) ** 2))
        C3_4_mean = np.mean(C3_4_array)
        C3_4_err = np.sqrt((nev - 1.0) / nev * np.sum((C3_4_array - C3_4_mean) ** 2))

        # Calculate vn{4} for each order and add small cutoff for security
        v1_4, v1_4_err = 0.0, 0.0
        if C1_4_mean < 0:
            v1_4 = (-C1_4_mean) ** 0.25 + 1e-30
            v1_4_err = 0.25 * (-C1_4_mean) ** -0.75 * C1_4_err

        v2_4, v2_4_err = 0.0, 0.0
        if C2_4_mean < 0:
            v2_4 = (-C2_4_mean) ** 0.25 + 1e-30
            v2_4_err = 0.25 * (-C2_4_mean) ** -0.75 * C2_4_err

        v3_4, v3_4_err = 0.0, 0.0
        if C3_4_mean < 0:
            v3_4 = (-C3_4_mean) ** 0.25 + 1e-30
            v3_4_err = 0.25 * (-C3_4_mean) ** -0.75 * C3_4_err

        results = [
            v1_4, v1_4_err, C1_4_mean, C1_4_err,
            v2_4, v2_4_err, C2_4_mean, C2_4_err,
            v3_4, v3_4_err, C3_4_mean, C3_4_err
        ]

        return results

    def calculate_vnSP_diff(self, QnpT_diff, Qnref):
        """
            this funciton calculates the scalar-product vn
            assumption: no overlap between particles of interest
                        and reference flow Qn vectors
            inputs: QnpT_diff[nev, norder, npT], Qnref[nev, norder]
            return: [vn{SP}(pT), vn{SP}(pT)_err]
        """
        QnpT_diff = np.array(QnpT_diff)
        Qnref = np.array(Qnref)
        nev, norder, npT = QnpT_diff.shape

        vn_values = []
        vn_errors = []

        Nref = np.real(Qnref[:, 0])
        N2refPairs = Nref*(Nref - 1.)
        NpTPOI = np.real(QnpT_diff[:, 0, :])
        N2POIPairs = NpTPOI*Nref.reshape(nev, 1)
        for iorder in range(1, norder):
            # compute Cn^ref{2}
            QnRef_tmp = Qnref[:, iorder]
            n2ref = np.abs(QnRef_tmp)**2. - Nref

            # compute vn{SP}(pT)
            QnpT_tmp = QnpT_diff[:, iorder, :]
            n2pT = np.real(QnpT_tmp*np.conj(QnRef_tmp.reshape(nev, 1)))

            # calcualte observables with Jackknife resampling method
            vnSPpT_arr = np.zeros([nev, npT])
            for iev in range(nev):
                array_idx = [True]*nev
                array_idx[iev] = False
                array_idx = np.array(array_idx)

                Cn2ref_arr = np.mean(n2ref[array_idx])/np.mean(N2refPairs[array_idx])
                vnSPpT_arr[iev, :] = (np.mean(n2pT[array_idx], 0)
                        /(np.mean(N2POIPairs[array_idx], 0)+1.e-20)/(np.sqrt(Cn2ref_arr))+1.e-20)
            vnSPpT_mean = np.mean(vnSPpT_arr, 0)
            vnSPpT_err  = np.sqrt((nev - 1.)/nev
                               *np.sum((vnSPpT_arr - vnSPpT_mean)**2., 0))
            
            vn_values.append(vnSPpT_mean)
            vn_errors.append(vnSPpT_err)

        return [vn_values, vn_errors]

    def calculate_vn4_diff(self, QnpT_diff, Qnref):
        """
            This function calculates the 4-particle vn(pT) using the scalar-product method.
            Assumption: No overlap between particles of interest and reference flow Qn vectors.
            Inputs:
                - QnpT_diff: Shape [nev, norder, npT]
                - Qnref: Shape [nev, norder]
            Returns:
                - vn_values: List of mean vn values for different harmonic orders (one array per harmonic order)
                - vn_errors: List of vn errors for different harmonic orders (one array per harmonic order)
        """
        QnpT_diff = np.array(QnpT_diff)
        Qnref = np.array(Qnref)
        nev, norder, npT = QnpT_diff.shape

        vn_values = []
        vn_errors = []

        for iorder in range(1, 4):  # Process for orders 1 to 3
            # compute Cn^ref{4}
            Nref = np.real(Qnref[:, 0])
            QnRef_tmp = Qnref[:, iorder]
            Q2nRef_tmp = Qnref[:, 2 * iorder]
            N4refPairs = Nref * (Nref - 1.) * (Nref - 2.) * (Nref - 3.)
            n4ref = (np.abs(QnRef_tmp)**4.
                     - 2. * np.real(Q2nRef_tmp * np.conj(QnRef_tmp) * np.conj(QnRef_tmp))
                     - 4. * (Nref - 2) * np.abs(QnRef_tmp)**2. + np.abs(Q2nRef_tmp)**2.
                     + 2. * Nref * (Nref - 3))
            N2refPairs = Nref * (Nref - 1.)
            n2ref = np.abs(QnRef_tmp)**2. - Nref

            # compute dn{4}(pT)
            NpTPOI = np.real(QnpT_diff[:, 0, :])
            QnpT_tmp = QnpT_diff[:, iorder, :]
            Nref = Nref.reshape(len(Nref), 1)
            QnRef_tmp = QnRef_tmp.reshape(len(QnRef_tmp), 1)
            Q2nRef_tmp = Q2nRef_tmp.reshape(len(Q2nRef_tmp), 1)
            N4POIPairs = NpTPOI * (Nref - 1.) * (Nref - 2.) * (Nref - 3.) + 1e-30
            n4pT = np.real(QnpT_tmp * QnRef_tmp * np.conj(QnRef_tmp) * np.conj(QnRef_tmp)
                           - 2. * (Nref - 1) * QnpT_tmp * np.conj(QnRef_tmp)
                           - QnpT_tmp * QnRef_tmp * np.conj(Q2nRef_tmp))
            N2POIPairs = NpTPOI * Nref + 1e-30
            n2pT = np.real(QnpT_tmp * np.conj(QnRef_tmp))

            # Calculate observables with Jackknife resampling
            Cn2ref_arr = np.zeros(nev)
            Cn4ref_arr = np.zeros(nev)
            dn4pT_arr = np.zeros(npT)
            vn4pT4_arr = np.zeros([nev, npT])
            
            for iev in range(nev):
                array_idx = [True] * nev
                array_idx[iev] = False
                array_idx = np.array(array_idx)

                Cn2ref_arr[iev] = (
                    np.mean(n2ref[array_idx]) / np.mean(N2refPairs[array_idx])
                )
                Cn4ref_arr[iev] = (
                    np.mean(n4ref[array_idx]) / np.mean(N4refPairs[array_idx])
                    - 2. * (Cn2ref_arr[iev])**2.
                )

                dn4pT_arr = (
                    np.mean(n4pT[array_idx, :], 0) / np.mean(N4POIPairs[array_idx, :], 0)
                    - 2. * np.mean(n2pT[array_idx, :], 0) / np.mean(N2POIPairs[array_idx, :], 0) * Cn2ref_arr[iev]
                )

                vn4pT4_arr[iev, :] = (-dn4pT_arr)**4. / ((-Cn4ref_arr[iev])**3.)

            vn4pT4_mean = np.mean(vn4pT4_arr, axis=0)
            vn4pT4_err = np.sqrt((nev - 1.) / nev * np.sum((vn4pT4_arr - vn4pT4_mean)**2., axis=0))

            vn4pT = np.zeros(npT)
            vn4pT_err = np.zeros(npT)
            idx = vn4pT4_mean > 0
            vn4pT[idx] = vn4pT4_mean[idx]**(0.25)
            vn4pT_err[idx] = vn4pT4_err[idx] / (4. * vn4pT4_mean[idx]**(0.75))

            # Append calculated values and errors for this harmonic order
            vn_values.append(vn4pT)
            vn_errors.append(vn4pT_err)

        return [vn_values, vn_errors]

    #-------------------------------------------------------------------------------------------
    # Plot distributions
    #-------------------------------------------------------------------------------------------
    def plot_observable(self, observable_type, observable, centrality, method='', pt_suffix='', logy = False):

        label = f'{observable_type}_{observable}{method}{self.suffix}_{centrality}{pt_suffix}'

        if observable_type == "soft_integrated":

            self.plot_distribution_and_ratio(observable_type, observable, centrality, logy)

        if observable_type == 'soft_differential':

            self.plot_distribution_and_ratio(observable_type, observable, centrality, logy)

        # # If AA: Plot PbPb/pp ratio, and comparison to data
        # # If pp: Plot distribution, and ratio to data
        # if self.is_AA:
        #     self.plot_RAA(observable_type, observable, centrality, label, pt_suffix=pt_suffix, logy=logy)
        #     self.write_experimental_data(observable_type, observable, centrality, label, pt_suffix=pt_suffix)
        # else:
        #     if self.observable_settings[f'jetscape_distribution']:
        #         self.plot_distribution_and_ratio(observable_type, observable, centrality, label, pt_suffix=pt_suffix, logy=logy)

    #-------------------------------------------------------------------------------------------
    # Write experimental data tables
    #-------------------------------------------------------------------------------------------
    def write_experimental_data(self, observable_type, observable, centrality, label, method='', pt_suffix=''):

        output_dir = os.path.join(self.output_dir, f'../Data')
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        force_write = True        
        filename = f'Data_{observable_type}_{observable}{method}{self.suffix}_{centrality}{pt_suffix}.dat'
        outputfile = os.path.join(output_dir, filename)
        if force_write or not os.path.exists(outputfile):

            # Get histogram binning
            key = [key for key in self.observable_settings.keys() if 'jetscape_distribution' in key and 'holes' not in key][0]
            h_prediction = self.observable_settings[key]

            # Truncate data tgraph to prediction histogram binning
            if h_prediction:

                xbins = np.array(h_prediction.GetXaxis().GetXbins())
                x_min = xbins[:-1]
                x_max = xbins[1:]
                
                g_data = self.observable_settings['data_distribution']
                g_truncated = self.plot_utils.truncate_tgraph(g_data, h_prediction, is_AA = self.is_AA) 
                if g_truncated:
                    y = np.array(g_truncated.GetY())
                    y_err = np.array([g_truncated.GetErrorY(i) for i in range(g_truncated.GetN())])
                    df = pd.DataFrame({'x_min': x_min, 'x_max': x_max, 'y': y, 'y_err': y_err})

                    x = np.array(g_truncated.GetX())
                    if np.any(np.greater(x_min, x)) or np.any(np.greater(x, x_max)):
                        sys.exit(f'ERROR: x not contained in hist binning: x={x} vs. x_bins={xbins}')

                    # Write table
                    header = f'Version 1.1\n'
                    header += f'Label xmin xmax y y_err'
                    np.savetxt(outputfile, df.values, header=header)
                else:
                    print(f'Did not write Data table: {filename} -- missing tgraph')
            else:
                print(f'Did not write Data table: {filename} -- missing histogram')
        else:
            sys.exit()

    #-------------------------------------------------------------------------------------------
    # Plot distributions in upper panel, and ratio in lower panel
    #-------------------------------------------------------------------------------------------
    def rebin_vn_results(self, pT_original, vn_values, vn_errors, pT_bins):
        """
        Rebin vn values and errors to match the provided pT_bins.
        
        Parameters:
        - pT_original: Original pT array (1D array).
        - vn_values: vn values (list of arrays, one for each harmonic order).
        - vn_errors: vn errors (list of arrays, one for each harmonic order).
        - pT_bins: The desired pT bins for rebinning (1D array).

        Returns:
        - rebinned_vn_values: List of rebinned vn values (one array per harmonic order).
        - rebinned_vn_errors: List of rebinned vn errors (one array per harmonic order).
        """
        rebinned_vn_values = []
        rebinned_vn_errors = []

        for vn_vals, vn_errs in zip(vn_values, vn_errors):
            rebinned_vals = []
            rebinned_errs = []

            for i in range(len(pT_bins) - 1):
                mask = (pT_original >= pT_bins[i]) & (pT_original < pT_bins[i + 1])
                if np.any(mask):
                    weights = 1 / (vn_errs[mask] ** 2 + 1e-20)  # Weight by inverse variance
                    mean_vn = np.sum(vn_vals[mask] * weights) / np.sum(weights)
                    err_vn = np.sqrt(1 / np.sum(weights))
                else:
                    mean_vn, err_vn = 0, 0

                rebinned_vals.append(mean_vn)
                rebinned_errs.append(err_vn)

            rebinned_vn_values.append(np.array(rebinned_vals))
            rebinned_vn_errors.append(np.array(rebinned_errs))

        return rebinned_vn_values, rebinned_vn_errors

    def plot_distribution_and_ratio(self, observable_type, observable, centrality, logy=False):
        """
        Plot data and model results for both soft_integrated (centrality-based) and soft_differential (pT-differential).
        Includes rebinning for pT-differential observables based on experimental pT bins.
        """
        # Determine the base name for fetching model results
        base_name = f"jetscape_distribution_{observable_type}_{observable}_{centrality}"
        print("base_name", base_name)
        if base_name not in self.observable_settings:
            print(f"WARNING: Model results for {base_name} not found.")
            return

        # Determine if observable is soft_integrated or soft_differential
        is_integrated = observable_type == "soft_integrated"

        # Fetch the experimental data bins
        if isinstance(self.observable_settings['data_distribution'], ROOT.TGraph):
            graph = self.observable_settings['data_distribution']
            n_points = graph.GetN()
            x_values = sorted([graph.GetX()[i] for i in range(n_points)])  # Sort x_values in ascending order

            if not is_integrated:
                # Directly compute bin edges using midpoints
                bin_edges = [x_values[0] - (x_values[1] - x_values[0]) / 2]  # First edge
                bin_edges.extend([(x_values[i - 1] + x_values[i]) / 2 for i in range(1, len(x_values))])  # Midpoints
                bin_edges.append(x_values[-1] + (x_values[-1] - x_values[-2]) / 2)  # Last edge
                pT_bins_data = np.array(bin_edges)
                print("Derived Bin Edges (pT_bins_data):", pT_bins_data)
                print("Bin Widths (Derived):", np.diff(pT_bins_data))
            else:
                # For integrated results, create centrality bin edges
                bin_width = x_values[-1] - x_values[-2] if len(x_values) > 1 else 1
                self.bins = np.array(x_values + [x_values[-1] + bin_width])
        else:
            print(f"WARNING: Data distribution not available or invalid for {observable}.")
            return

        # Fetch the model results
        model_results = self.observable_settings[base_name]

        # Process soft_integrated observables (no rebinning required)
        if is_integrated:
            # Centrality midpoint and error
            centrality_midpoint = (centrality[0] + centrality[1]) / 2.0
            centrality_width = (centrality[1] - centrality[0]) / 2.0

            # Determine the y-value and error based on the observable
            if 'multiplicity' in observable:
                y_value = model_results
                y_error = 0  # Assuming no error for multiplicity
            elif 'v2' in observable:
                if 'two' in observable:
                    vn2, vn2_err = model_results
                    y_value = vn2[1]  # vn_2 for n=2
                    y_error = vn2_err[1]
                elif 'four' in observable:
                    vn4 = model_results
                    y_value = vn4[4]  # v2_4 for n=2
                    y_error = vn4[5]  # v2_4_err for n=2
            else:
                print(f"WARNING: Unrecognized soft_integrated observable {observable}.")
                return

            # Create a TGraphErrors for centrality-based results
            model_graph = ROOT.TGraphErrors(1)
            model_graph.SetPoint(0, centrality_midpoint, y_value)
            model_graph.SetPointError(0, centrality_width, y_error)

        # Process soft_differential observables (with rebinning)
        else:
            if 'v2' in observable:
                pT_model = model_results['pT_bins']
                vn_values = model_results['vn_values']
                vn_errors = model_results['vn_errors']

                # Rebin the model results
                vn_values_rebinned, vn_errors_rebinned = self.rebin_vn_results(
                    pT_model, vn_values, vn_errors, pT_bins_data
                )

                model_graph = ROOT.TGraphErrors(len(x_values))
                for i, (vn, vn_err) in enumerate(zip(vn_values_rebinned[1], vn_errors_rebinned[1])):  # n=2
                    bin_center = x_values[i]  # Use experimental bin centers
                    bin_width = (pT_bins_data[i + 1] - pT_bins_data[i]) / 2  # Use derived edges for widths
                    model_graph.SetPoint(i, bin_center, vn)
                    model_graph.SetPointError(i, bin_width, vn_err)

        # Style the model graph
        model_graph.SetMarkerSize(self.marker_size)
        model_graph.SetMarkerStyle(self.data_marker + 1)
        model_graph.SetMarkerColor(ROOT.kRed)
        model_graph.SetLineStyle(self.line_style)
        model_graph.SetLineWidth(self.line_width)

        # Create a canvas for the plot
        c = ROOT.TCanvas('c', 'c', 1000, 950)
        c.Draw()
        c.cd()

        # Configure the top pad
        pad1 = ROOT.TPad('myPad', 'The pad', 0, 0, 1, 1)
        pad1.SetLeftMargin(0.2)
        pad1.SetTopMargin(0.08)
        pad1.SetRightMargin(0.04)
        pad1.SetBottomMargin(0.15)
        pad1.SetTicks(0, 1)
        pad1.Draw()
        if logy:
            pad1.SetLogy()
        pad1.cd()

        # Configure and draw the blank histogram for the axes
        if is_integrated:
            x_min, x_max = self.bins[0], self.bins[-1]  # Extend centrality range slightly for visualization
        else:
            x_min, x_max = pT_bins_data[0], pT_bins_data[-1]

        myBlankHisto = ROOT.TH1F('myBlankHisto', 'Blank Histogram', 1, x_min, x_max)
        myBlankHisto.SetNdivisions(505)
        myBlankHisto.SetXTitle(self.xtitle)
        myBlankHisto.SetYTitle(self.ytitle)
        myBlankHisto.SetMaximum(self.y_max)
        myBlankHisto.SetMinimum(self.y_min)
        myBlankHisto.GetYaxis().SetRangeUser(self.y_min, self.y_max)
        myBlankHisto.GetYaxis().SetTitleSize(0.08)
        myBlankHisto.GetYaxis().SetTitleOffset(1.1)
        myBlankHisto.GetYaxis().SetLabelSize(0.06)
        myBlankHisto.GetXaxis().SetTitleSize(0.08)
        myBlankHisto.GetXaxis().SetTitleOffset(1.1)
        myBlankHisto.GetXaxis().SetLabelSize(0.06)
        myBlankHisto.Draw('E')

        # Draw the data distribution
        self.observable_settings['data_distribution'].SetMarkerSize(self.marker_size)
        self.observable_settings['data_distribution'].SetMarkerStyle(self.data_marker)
        self.observable_settings['data_distribution'].SetMarkerColor(self.data_color)
        self.observable_settings['data_distribution'].SetLineStyle(self.line_style)
        self.observable_settings['data_distribution'].SetLineWidth(self.line_width)
        self.observable_settings['data_distribution'].SetLineColor(self.data_color)
        self.observable_settings['data_distribution'].Draw('PE Z same')

        # Draw the model graph
        model_graph.Draw('PE Z same')

        # Add legend
        legend = ROOT.TLegend(0.4, 0.75, 0.75, 0.85)
        self.plot_utils.setup_legend(legend, 0.045, sep=0.001)
        legend.AddEntry(self.observable_settings['data_distribution'], 'Data', 'PE')
        legend.AddEntry(model_graph, 'Model', 'PE')
        legend.Draw()

        # Add text annotations
        pad1.cd()
        text_latex = ROOT.TLatex()
        text_latex.SetNDC()
        text_latex.SetTextSize(0.065)
        text = f'#bf{{{observable_type}_{observable}}} #sqrt{{#it{{s}}}} = {self.sqrts / 1000.} TeV'
        text_latex.DrawLatex(0.25, 0.83, text)
        text = f'{centrality}'
        text_latex.DrawLatex(0.25, 0.73, text)

        # Save the canvas
        hname = f'h_{observable_type}_{observable}_{centrality}'
        c.SaveAs(os.path.join(self.output_dir, f'{hname}{self.file_format}'))
        c.Close()

    #-------------------------------------------------------------------------------------------
    # Plot event QA
    #-------------------------------------------------------------------------------------------
    def plot_event_qa(self):

        if self.is_AA:

            for centrality in self.observable_centrality_list:

                # Only plot those centralities that exist
                if np.isclose(self.input_file.Get(f'h_centrality_generated').Integral(centrality[0]+1, centrality[1]), 0):
                    continue
                print(centrality)

                # Crosscheck that pt-hat weighting is satisfactory
                # Make sure that the large-weight, low-pt-hat range has sufficient statistics to avoid normalization fluctuations
                sum_weights = self.input_file.Get(f'h_weight_sum_{centrality}').GetBinContent(1)
                h_pt_hat = self.input_file.Get('h_pt_hat')
                h_pt_hat_weighted = self.input_file.Get('h_pt_hat_weighted')

                # Normalize by sum of weights, i.e. 1/sigma_pt_hat * dsigma/dpt_hat
                h_pt_hat.Scale(1./sum_weights, 'width')
                h_pt_hat_weighted.Scale(1./sum_weights, 'width')

                # Compute normalization uncertainty (binned approximation)
                h_weights = self.input_file.Get(f'h_weights_{centrality}')
                sum_weights_integral = 0
                normalization_uncertainty = 0
                for i in range(1, h_weights.GetNbinsX()+1):
                    sum_weights_integral += h_weights.GetBinCenter(i)*h_weights.GetBinContent(i)
                    normalization_uncertainty = np.sqrt( np.square(normalization_uncertainty) + h_weights.GetBinContent(i)*np.square(h_weights.GetBinCenter(i)) )
                print(f'sum_weights {centrality}: {sum_weights}')
                print(f'sum_weights_integral {centrality}: {sum_weights_integral}')
                print(f'normalization_uncertainty {centrality}: {100*normalization_uncertainty/sum_weights_integral} %')

                # Also compute overall pt-hat cross-section uncertainty
                h_xsec = self.input_file.Get(f'h_xsec_{centrality}')
                xsec = h_xsec.GetBinContent(1)
                xsec_error = self.input_file.Get(f'h_xsec_error_{centrality}').GetBinContent(1)
                print(f'xsec {centrality}: {xsec/h_xsec.GetEntries()}')
                print(f'xsec_uncertainty {centrality}: {100*xsec_error/xsec}')
                print()

        else:

            # Crosscheck that pt-hat weighting is satisfactory
            # Make sure that the large-weight, low-pt-hat range has sufficient statistics to avoid normalization fluctuations
            sum_weights = self.input_file.Get('h_weight_sum').GetBinContent(1)
            h_pt_hat = self.input_file.Get('h_pt_hat')
            h_pt_hat_weighted = self.input_file.Get('h_pt_hat_weighted')

            # Normalize by sum of weights, i.e. 1/sigma_pt_hat * dsigma/dpt_hat
            h_pt_hat.Scale(1./sum_weights, 'width')
            h_pt_hat_weighted.Scale(1./sum_weights, 'width')

            # Compute normalization uncertainty (binned approximation)
            h_weights = self.input_file.Get('h_weights')
            sum_weights_integral = 0
            normalization_uncertainty = 0
            for i in range(1, h_weights.GetNbinsX()+1):
                sum_weights_integral += h_weights.GetBinCenter(i)*h_weights.GetBinContent(i)
                normalization_uncertainty = np.sqrt( np.square(normalization_uncertainty) + h_weights.GetBinContent(i)*np.square(h_weights.GetBinCenter(i)) )
            print(f'sum_weights: {sum_weights}')
            print(f'sum_weights_integral: {sum_weights_integral}')
            print(f'normalization_uncertainty: {100*normalization_uncertainty/sum_weights_integral} %')

            # Also compute overall pt-hat cross-section uncertainty
            h_xsec = self.input_file.Get('h_xsec')
            xsec = h_xsec.GetBinContent(1)
            xsec_error = self.input_file.Get('h_xsec_error').GetBinContent(1)
            print(f'xsec: {xsec/h_xsec.GetEntries()}')
            print(f'xsec_uncertainty: {100*xsec_error/xsec}')

        c = ROOT.TCanvas('c', 'c', 600, 650)
        c.Draw()
        c.cd()

        # Distribution
        pad2_dy = 0.45
        pad1 = ROOT.TPad('myPad', 'The pad',0,pad2_dy,1,1)
        pad1.SetLeftMargin(0.2)
        pad1.SetTopMargin(0.08)
        pad1.SetRightMargin(0.04)
        pad1.SetBottomMargin(0.)
        pad1.SetTicks(0,1)
        pad1.Draw()
        pad1.SetLogy()
        pad1.SetLogx()
        pad1.cd()

        legend = ROOT.TLegend(0.58,0.65,0.75,0.9)
        self.plot_utils.setup_legend(legend, 0.055, sep=-0.1, title=f'#sqrt{{s_{{NN}}}} = {self.sqrts} GeV')

        self.bins = np.array(h_pt_hat.GetXaxis().GetXbins())
        myBlankHisto = ROOT.TH1F('myBlankHisto','Blank Histogram', 1, self.bins[0], self.bins[-1])
        myBlankHisto.SetNdivisions(505)
        myBlankHisto.SetXTitle('#hat{p}_{T} (GeV/#it{c})')
        myBlankHisto.SetYTitle('#frac{1}{#sigma_{#hat{p}_{T}}} #frac{d#sigma}{d#hat{p}_{T}}')
        myBlankHisto.SetMaximum(np.power(10,5+(self.power-3)))
        myBlankHisto.SetMinimum(2e-15) # Don't draw 0 on top panel
        myBlankHisto.GetYaxis().SetTitleSize(0.08)
        myBlankHisto.GetYaxis().SetTitleOffset(1.1)
        myBlankHisto.GetYaxis().SetLabelSize(0.06)
        myBlankHisto.Draw('E')

        # Ratio
        c.cd()
        pad2 = ROOT.TPad('pad2', 'pad2', 0, 0.02, 1, pad2_dy)
        pad2.SetTopMargin(0)
        pad2.SetBottomMargin(0.21)
        pad2.SetLeftMargin(0.2)
        pad2.SetRightMargin(0.04)
        pad2.SetTicks(0,1)
        pad2.SetLogx()
        pad2.Draw()
        pad2.cd()

        myBlankHisto2 = myBlankHisto.Clone('myBlankHisto_C')
        myBlankHisto2.SetYTitle('#frac{weighted}{ideal}')
        myBlankHisto2.SetXTitle('#hat{p}_{T} (GeV/#it{c})')
        myBlankHisto2.GetXaxis().SetTitleSize(26)
        myBlankHisto2.GetXaxis().SetTitleFont(43)
        myBlankHisto2.GetXaxis().SetTitleOffset(2.3)
        myBlankHisto2.GetXaxis().SetLabelFont(43)
        myBlankHisto2.GetXaxis().SetLabelSize(22)
        myBlankHisto2.GetYaxis().SetTitleSize(28)
        myBlankHisto2.GetYaxis().SetTitleFont(43)
        myBlankHisto2.GetYaxis().SetTitleOffset(2.)
        myBlankHisto2.GetYaxis().SetLabelFont(43)
        myBlankHisto2.GetYaxis().SetLabelSize(20)
        myBlankHisto2.GetYaxis().SetNdivisions(505)
        myBlankHisto2.GetYaxis().SetRangeUser(0., 1.99)
        myBlankHisto2.Draw('')

        # Draw generated pt-hat
        pad1.cd()
        h_pt_hat.SetMarkerSize(self.marker_size)
        h_pt_hat.SetMarkerStyle(self.data_marker)
        h_pt_hat.SetMarkerColor(self.data_color)
        h_pt_hat.SetLineStyle(self.line_style)
        h_pt_hat.SetLineWidth(self.line_width)
        h_pt_hat.SetLineColor(self.data_color)
        h_pt_hat.Draw('PE Z same')
        legend.AddEntry(h_pt_hat, f'#alpha={self.power}, generated', 'PE')

        # Draw weighted pt-hat
        h_pt_hat_weighted.SetMarkerSize(self.marker_size)
        h_pt_hat_weighted.SetMarkerStyle(self.data_marker)
        h_pt_hat_weighted.SetMarkerColor(self.jetscape_color[0])
        h_pt_hat_weighted.SetLineStyle(self.line_style)
        h_pt_hat_weighted.SetLineWidth(self.line_width)
        h_pt_hat_weighted.SetLineColor(self.jetscape_color[0])
        h_pt_hat_weighted.Draw('PE Z same')
        legend.AddEntry(h_pt_hat_weighted, f'#alpha={self.power}, weighted', 'PE')

        # Draw generated pt-hat weighted by inverse of ideal weight (pthat/ptref)^alpha
        h_pt_hat_weighted_ideal = h_pt_hat.Clone()
        h_pt_hat_weighted_ideal.SetName('h_pt_hat_weighted_ideal')
        for i in range(1, h_pt_hat_weighted_ideal.GetNbinsX()+1):
            content = h_pt_hat_weighted_ideal.GetBinContent(i)
            low_edge = h_pt_hat_weighted_ideal.GetXaxis().GetBinLowEdge(i)
            up_edge = h_pt_hat_weighted_ideal.GetXaxis().GetBinUpEdge(i)
            ideal_weight = ( np.power(up_edge, self.power+1)-np.power(low_edge, self.power+1) ) / ( (up_edge-low_edge)*np.power(self.pt_ref, self.power)*(self.power+1) )
            #ideal_weight = np.power( h_pt_hat_weighted_ideal.GetXaxis().GetBinCenter(i) / self.pt_ref, self.power)
            h_pt_hat_weighted_ideal.SetBinContent(i, content / ideal_weight)
            h_pt_hat_weighted_ideal.SetBinError(i, 0.)
        h_pt_hat_weighted_ideal.SetLineColorAlpha(ROOT.kGreen-8, self.jetscape_alpha[0])
        h_pt_hat_weighted_ideal.SetLineWidth(3)
        h_pt_hat_weighted_ideal.Draw('L same')
        legend.AddEntry(h_pt_hat_weighted_ideal, f'#alpha={self.power}, ideal weight', 'L')

        legend.Draw()

        # Plot ratio of weighted to ideal weighted
        pad2.cd()
        h_pt_hat_weighted_ratio = h_pt_hat_weighted.Clone()
        h_pt_hat_weighted_ratio.SetName('h_pt_hat_weighted_ratio')
        h_pt_hat_weighted_ratio.Divide(h_pt_hat_weighted_ideal)
        h_pt_hat_weighted_ratio.SetMarkerSize(self.marker_size)
        h_pt_hat_weighted_ratio.SetMarkerStyle(self.data_marker)
        h_pt_hat_weighted_ratio.SetMarkerColor(self.jetscape_color[0])
        h_pt_hat_weighted_ratio.SetLineStyle(self.line_style)
        h_pt_hat_weighted_ratio.SetLineWidth(self.line_width)
        h_pt_hat_weighted_ratio.SetLineColor(self.jetscape_color[0])
        h_pt_hat_weighted_ratio.Draw('PE Z same')

        line = ROOT.TLine(self.bins[0], 1, self.bins[-1], 1)
        line.SetLineColor(920+2)
        line.SetLineStyle(2)
        line.SetLineWidth(2)
        line.Draw()

        pad1.cd()
        text_latex = ROOT.TLatex()
        text_latex.SetNDC()

        x = 0.25
        text_latex.SetTextSize(0.06)
        text = '#sigma_{{n_{{event}}}} = {:.2f}%'.format(100*normalization_uncertainty/sum_weights_integral)
        text_latex.DrawLatex(x, 0.83, text)
        text = '#sigma_{{xsec}} = {:.2f}%'.format(100*xsec_error/xsec)
        text_latex.DrawLatex(x, 0.76, text)

        c.SaveAs(os.path.join(self.output_dir, f'pt_hat{self.file_format}'))
        c.Close()

    # ---------------------------------------------------------------
    # Save all ROOT histograms to file
    # ---------------------------------------------------------------
    def write_output_objects(self):
        print()

        # Save histograms to ROOT file
        output_ROOT_filename = os.path.join(self.output_dir, 'final_results.root')
        f_ROOT = ROOT.TFile(output_ROOT_filename, 'recreate')
        f_ROOT.cd()
        for key,val in self.output_dict.items():
            if val:
                val.SetName(key)
                val.Write()
                if isinstance(val, (ROOT.TH1)):
                    val.SetDirectory(0)
                    del val
        f_ROOT.Close()

        # Also save JETSCAPE AA/pp ratios as numpy arrays to HDF5 file
        if self.is_AA:
            with uproot.open(output_ROOT_filename) as uproot_file:
                output_HDF5_filename = os.path.join(self.output_dir, 'final_results.h5')
                with h5py.File(output_HDF5_filename, 'w') as hf:

                    keys = [key.split(';',1)[0] for key in uproot_file.keys()]
                    for key in keys:
                        if 'jetscape' in key:

                            bin_edges = uproot_file[key].axis().edges()
                            values = uproot_file[key].values()
                            errors = uproot_file[key].errors()

                            hf.create_dataset(f'{key}_bin_edges', data=bin_edges)
                            hf.create_dataset(f'{key}_values', data=values)
                            hf.create_dataset(f'{key}_errors', data=errors)

    #-------------------------------------------------------------------------------------------
    # Generate pptx of one plot per slide, for convenience
    #-------------------------------------------------------------------------------------------
    def generate_pptx(self):
        # Delayed import since this isn't on available on many systems
        # For those that need can install it (ie. systems where powerpoint can run) you can use:
        #   pip install python-pptx
        import pptx

        # Create a blank presentation
        p = pptx.Presentation()

        # Set slide layouts
        title_slide_layout = p.slide_layouts[0]
        blank_slide_layout = p.slide_layouts[6]

        # Make a title slide
        slide = p.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = f'QA for {self.sqrts/1000.} TeV'
        author = slide.placeholders[1]
        author.text = 'STAT WG'

        # Loop through all output files and plot
        files = [f for f in os.listdir(self.output_dir) if f.endswith(self.file_format)]
        for file in sorted(files):
            img = os.path.join(self.output_dir, file)
            slide = p.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img, left=pptx.util.Inches(2.),
                                          top=pptx.util.Inches(1.),
                                          width=pptx.util.Inches(5.))

        p.save(os.path.join(self.output_dir, 'results.pptx'))

#-------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------
if __name__ == '__main__':
    print('Executing plot_results_STAT.py...')
    print('')

    # Define arguments
    parser = argparse.ArgumentParser(description='Plot JETSCAPE events')
    parser.add_argument(
        '-c',
        '--configFile',
        action='store',
        type=str,
        metavar='configFile',
        default='config/TG3.yaml',
        help='Config file'
    )
    parser.add_argument(
        '-i',
        '--inputFile',
        action='store',
        type=str,
        metavar='inputFile',
        default='pp_5020_plot/histograms_5020_merged.root',
        help='Input file'
    )
    parser.add_argument(
        '-o',
        '--outputDir',
        action='store',
        type=str,
        metavar='outputDir',
        default='',
        help='Output directory for output to be written to',
    )
    parser.add_argument(
        '-r',
        '--refFile',
        action='store',
        type=str,
        metavar='refFile',
        default='final_results.root',
        help='pp reference file'
    )

    # Parse the arguments
    args = parser.parse_args()

    # If invalid configFile is given, exit
    if not os.path.exists(args.configFile):
        print('File "{0}" does not exist! Exiting!'.format(args.configFile))
        sys.exit(0)

    # If invalid inputDir is given, exit
    if not os.path.exists(args.inputFile):
        print('File "{0}" does not exist! Exiting!'.format(args.inputFile))
        sys.exit(0)

    analysis = PlotResults(config_file=args.configFile, input_file=args.inputFile, output_dir=args.outputDir)
    analysis.plot_results()
